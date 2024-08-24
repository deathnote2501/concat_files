import streamlit as st
import os
from io import StringIO
import pdfplumber
from docx import Document
from pptx import Presentation
import tempfile
import base64

# Define the password
PASSWORD = st.secrets["PASSWORD"]

def read_pdf(file):
    try:
        content = ""
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                content += page.extract_text()
        return content
    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier PDF : {str(e)}")
        return ""

def read_docx(file):
    try:
        doc = Document(file)
        content = "\n".join([para.text for para in doc.paragraphs])
        return content
    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier DOCX : {str(e)}")
        return ""

def read_txt(file):
    try:
        return file.read().decode('utf-8')
    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier texte : {str(e)}")
        return ""

def read_pptx(file):
    try:
        prs = Presentation(file)
        content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        return "\n".join(content)
    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier PPTX : {str(e)}")
        return ""

def save_concatenated_file(content):
    try:
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
        with open(temp_file.name, 'w', encoding='utf-8') as f:
            f.write(content)
        return temp_file.name
    except Exception as e:
        st.error(f"Erreur lors de l'enregistrement du fichier : {str(e)}")
        return None

def main():
    st.markdown("<h1 style='text-align: center;'>Création base de connaissances pour les GPTs</h1>", unsafe_allow_html=True)
    st.markdown("<p style='text-align: center;'>Par Jérome IAvarone - IAvaronce conseil</p>", unsafe_allow_html=True)
    st.write("")

    image_url = "https://www.iacademy-formation.com/wp-content/uploads/2024/08/iyus-sugiharto-jpkxJAcp6a4-unsplash-modified-1.png"
    st.image(image_url, use_column_width=True)


    # Password input
    password = st.text_input("Entrez le mot de passe :", type="password")

    if password == PASSWORD:
        st.write("")
        st.write("")
        st.write("")
        st.write("")
        st.write("")
        st.markdown("<h2 style='text-align: left;'>Chargez vos fichiers PDF</h2>", unsafe_allow_html=True)


        uploaded_files = st.file_uploader("", type=["pdf", "doc", "docx", "txt", "ppt", "pptx"], accept_multiple_files=True)
        
        if st.button("Créer sa base de connaissances"):
            if not uploaded_files:
                st.warning("Chargez au moins 1 fichier.")
            else:
                all_content = ""
                for uploaded_file in uploaded_files:
                    file_name = uploaded_file.name
                    file_extension = os.path.splitext(file_name)[1].lower()
                    try:
                        if file_extension == ".pdf":
                            file_content = read_pdf(uploaded_file)
                        elif file_extension in [".doc", ".docx"]:
                            file_content = read_docx(uploaded_file)
                        elif file_extension == ".txt":
                            file_content = read_txt(uploaded_file)
                        elif file_extension in [".ppt", ".pptx"]:
                            file_content = read_pptx(uploaded_file)
                        else:
                            file_content = ""
                        
                        all_content += f"--------------------------- BEGIN {file_name} ---------------------------\n"
                        all_content += file_content
                        all_content += f"\n--------------------------- END {file_name} ---------------------------\n\n"
                    except Exception as e:
                        st.error(f"Erreur lors du traitement du fichier {file_name} : {str(e)}")
                
                concatenated_file_path = save_concatenated_file(all_content)
                
                if concatenated_file_path:
                    st.write("\n")
                    st.success("Traitement réalisé avec succès :)")
                    st.write("\n")

                    with open(concatenated_file_path, 'rb') as f:
                        b64 = base64.b64encode(f.read()).decode()
                        href = f'<a href="data:file/txt;base64,{b64}" download="concatenated_file.txt" style="font-size:20px;">>> Télécharger sa base de connaissances</a>'
                        st.markdown(href, unsafe_allow_html=True)

        st.write("")
        st.write("")
        st.write("")
        st.markdown("<p style='text-align: center;'>© 2024 Jérome IAvarone - jerome.iavarone@gmail.com</p>", unsafe_allow_html=True)
    elif password:
        st.error("Mot de passe incorrect")

if __name__ == "__main__":
    main()
